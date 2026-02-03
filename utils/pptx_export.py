# Fonction utilitaire pour exporter les réponses en présentation PowerPoint
import re
from io import BytesIO
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def detect_presentation_structure(text: str) -> Dict:
    """
    Détecte la structure d'une présentation dans le texte.

    Args:
        text: Texte à analyser

    Returns:
        Dict avec 'has_slides', 'has_bullets', 'has_tables', 'slide_count'
    """
    result = {
        'has_slides': False,
        'has_bullets': False,
        'has_tables': False,
        'slide_count': 0
    }

    # Détection de slides (markdown headers ou numérotation)
    slide_patterns = [
        r'^#{1,2}\s+',  # Headers markdown
        r'^Slide\s+\d+',  # "Slide 1", "Slide 2"
        r'^\d+\.\s+[A-Z]',  # "1. Titre", "2. Autre titre"
    ]

    lines = text.split('\n')
    for line in lines:
        for pattern in slide_patterns:
            if re.match(pattern, line, re.IGNORECASE | re.MULTILINE):
                result['has_slides'] = True
                result['slide_count'] += 1
                break

    # Détection de listes à puces
    if re.search(r'^\s*[-*•]\s+', text, re.MULTILINE):
        result['has_bullets'] = True

    # Détection de tableaux
    if '|' in text and text.count('|') > 3:
        result['has_tables'] = True

    return result


def extract_slides_content(text: str) -> List[Dict]:
    """
    Extrait le contenu de chaque slide du texte.

    Args:
        text: Texte contenant la structure de présentation

    Returns:
        Liste de dict avec 'title', 'content', 'bullets', 'table'
    """
    slides = []
    current_slide = None
    lines = text.split('\n')

    # Patterns pour identifier les titres de slides
    title_patterns = [
        (r'^##\s+(.+)$', 'markdown_h2'),  # ## Titre
        (r'^#\s+(.+)$', 'markdown_h1'),   # # Titre
        (r'^Slide\s+\d+\s*:\s*(.+)$', 'slide_numbered'),  # Slide 1: Titre
        (r'^\*\*(.+)\*\*$', 'bold_title'),  # **Titre**
    ]

    for line in lines:
        line_stripped = line.strip()

        # Ignorer les lignes vides
        if not line_stripped:
            if current_slide and current_slide['content']:
                current_slide['content'].append('')
            continue

        # Vérifier si c'est un titre de slide
        is_title = False
        for pattern, pattern_type in title_patterns:
            match = re.match(pattern, line_stripped)
            if match:
                # Sauvegarder le slide précédent
                if current_slide:
                    slides.append(current_slide)

                # Créer un nouveau slide
                title = match.group(1).strip()
                current_slide = {
                    'title': title,
                    'content': [],
                    'bullets': [],
                    'table': None
                }
                is_title = True
                break

        if is_title:
            continue

        # Si pas encore de slide, créer un slide par défaut
        if current_slide is None:
            current_slide = {
                'title': 'Slide 1',
                'content': [],
                'bullets': [],
                'table': None
            }

        # Détecter les puces
        bullet_match = re.match(r'^\s*[-*•]\s+(.+)$', line_stripped)
        if bullet_match:
            current_slide['bullets'].append(bullet_match.group(1))
        else:
            current_slide['content'].append(line_stripped)

    # Ajouter le dernier slide
    if current_slide:
        slides.append(current_slide)

    return slides


def extract_table_from_text(text: str) -> Optional[List[List[str]]]:
    """
    Extrait un tableau markdown du texte.

    Args:
        text: Texte contenant potentiellement un tableau

    Returns:
        Tableau sous forme de liste de listes ou None
    """
    lines = text.split('\n')
    table_data = []

    for line in lines:
        if '|' in line:
            # Ignorer les lignes de séparation
            if re.match(r'^\s*\|[\s\-:|]+\|\s*$', line):
                continue

            # Extraire les cellules
            cells = [cell.strip() for cell in line.split('|')]
            cells = [c for c in cells if c]

            if cells:
                table_data.append(cells)

    return table_data if table_data else None


# Crée une présentation PowerPoint à partir de la réponse.
# Args:
#     question: Question posée (non utilisée, conservée pour compatibilité)
#     response: Réponse générée contenant les slides
#     metadata: Métadonnées (non utilisées, conservées pour compatibilité)
# Returns:
#     Buffer contenant le fichier PowerPoint
def create_powerpoint_export(question: str, response: str, metadata: Dict) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Couleurs personnalisées
    BLUE = RGBColor(54, 96, 146)
    DARK_GRAY = RGBColor(68, 68, 68)

    # Détecter la structure
    structure = detect_presentation_structure(response)

    # Extraire et créer les slides de contenu UNIQUEMENT
    if structure['has_slides']:
        slides_content = extract_slides_content(response)

        for slide_data in slides_content:
            # Utiliser le layout vierge (6) pour éviter les placeholders par défaut
            slide_layout = prs.slide_layouts[6]  # Layout vide
            slide = prs.slides.add_slide(slide_layout)

            # Ajouter le titre manuellement
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = BLUE

            # Zone de contenu
            content_left = Inches(0.5)
            content_top = Inches(2)
            content_width = Inches(9)
            content_height = Inches(5)

            # Si des puces existent
            if slide_data['bullets']:
                content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                tf = content_box.text_frame
                tf.word_wrap = True

                for i, bullet_text in enumerate(slide_data['bullets']):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = bullet_text
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = DARK_GRAY
                    p.space_before = Pt(12)

            # Sinon, ajouter le contenu texte
            elif slide_data['content']:
                content_text = '\n'.join(slide_data['content'])

                # Vérifier s'il y a un tableau
                table_data = extract_table_from_text(content_text)

                if table_data:
                    # Créer un tableau
                    rows = len(table_data)
                    cols = len(table_data[0]) if table_data else 0

                    if rows > 0 and cols > 0:
                        table_left = Inches(1)
                        table_top = Inches(2.5)
                        table_width = Inches(8)
                        table_height = Inches(0.5) * rows

                        table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
                        table = table_shape.table

                        # Remplir le tableau
                        for i, row_data in enumerate(table_data):
                            for j, cell_text in enumerate(row_data):
                                if j < cols:  # Sécurité pour éviter les dépassements
                                    cell = table.cell(i, j)
                                    cell.text = cell_text

                                    # Style en-tête
                                    if i == 0:
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = BLUE
                                        paragraph = cell.text_frame.paragraphs[0]
                                        paragraph.font.bold = True
                                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                        paragraph.font.size = Pt(14)
                                    else:
                                        paragraph = cell.text_frame.paragraphs[0]
                                        paragraph.font.size = Pt(12)
                                        paragraph.font.color.rgb = DARK_GRAY

                else:
                    # Texte simple
                    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                    tf = content_box.text_frame
                    tf.text = content_text
                    tf.word_wrap = True

                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(16)
                        paragraph.font.color.rgb = DARK_GRAY
                        paragraph.space_before = Pt(6)

    else:
        # Pas de structure détectée - créer un slide unique avec tout le contenu
        slide_layout = prs.slide_layouts[6]  # Layout vide
        slide = prs.slides.add_slide(slide_layout)

        # Titre
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = Inches(9)
        title_height = Inches(1)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = "Réponse"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = BLUE

        # Contenu
        content_left = Inches(0.5)
        content_top = Inches(2)
        content_width = Inches(9)
        content_height = Inches(5)

        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf = content_box.text_frame
        tf.word_wrap = True

        # Diviser le texte en paragraphes
        paragraphs = response.split('\n\n')

        for i, para_text in enumerate(paragraphs[:10]):  # Limiter à 10 paragraphes
            if not para_text.strip():
                continue

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = para_text.strip()
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(12)

    # Sauvegarder dans un buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer
